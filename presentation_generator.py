#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентаций из структурированного текста
Поддерживает формат с маркерами ##-TOPIC-START-##, #-SLIDE-START-#
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import argparse
import os


class PresentationGenerator:
    def __init__(self):
        # Настройки стилей согласно требованиям
        self.title_font_size = Pt(28)
        self.subtitle_font_size = Pt(20)
        self.content_font_size = Pt(14)  # Требование: кегль 14 для основного текста
        self.code_font_size = Pt(12)
        
        # Цвета
        self.title_color = RGBColor(31, 73, 125)  # Темно-синий
        self.subtitle_color = RGBColor(68, 114, 196)  # Синий
        self.content_color = RGBColor(68, 68, 68)  # Темно-серый
        self.code_color = RGBColor(0, 100, 0)  # Зеленый для кода
        
        # Ключевые слова для выделения жирным
        self.bold_keywords = ['Цель:', 'Задачи:', 'Пример', 'Задание', 'Шаг']
        
    def parse_content(self, content):
        """Парсинг содержимого по новому формату"""
        presentations = []
        current_presentation = None
        current_slide = None
        
        lines = content.split('\n')
        
        for line in lines:
            line = line.strip()
            
            if line.startswith('##-TOPIC-START-##'):
                # Сохраняем предыдущую презентацию
                if current_presentation is not None:
                    if current_slide is not None:
                        current_presentation['slides'].append(current_slide)
                    presentations.append(current_presentation)
                
                # Начинаем новую презентацию
                current_presentation = {
                    'title': '',
                    'level': '',
                    'module': '',
                    'slides': []
                }
                current_slide = None
                
            elif line.startswith('#-SLIDE-START-#'):
                # Сохраняем предыдущий слайд
                if current_slide is not None and current_presentation is not None:
                    current_presentation['slides'].append(current_slide)
                
                # Начинаем новый слайд
                current_slide = {
                    'title': '',
                    'content': []
                }
                
            elif line.startswith('TITLE::'):
                # Заголовок слайда
                if current_slide is not None:
                    current_slide['title'] = line.replace('TITLE::', '').strip()
                    
            elif line.startswith('Практическая работа'):
                # Заголовок презентации
                if current_presentation is not None:
                    current_presentation['title'] = line
                    
            elif line.startswith('Уровень:'):
                # Уровень презентации
                if current_presentation is not None:
                    current_presentation['level'] = line.replace('Уровень:', '').strip()
                    
            elif line.startswith('Модуль'):
                # Модуль презентации
                if current_presentation is not None:
                    current_presentation['module'] = line
                    
            elif line and current_slide is not None:
                # Добавляем контент к текущему слайду
                current_slide['content'].append(line)
        
        # Сохраняем последние слайд и презентацию
        if current_slide is not None and current_presentation is not None:
            current_presentation['slides'].append(current_slide)
        if current_presentation is not None:
            presentations.append(current_presentation)
            
        return presentations
    
    def create_title_slide(self, prs, title, subtitle=""):
        """Создание титульного слайда"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.name = 'Calibri'
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = self.subtitle_font_size
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = self.subtitle_color
            subtitle_shape.text_frame.paragraphs[0].font.name = 'Calibri'
    
    def should_bold_text(self, text):
        """Проверка, нужно ли выделять текст жирным"""
        for keyword in self.bold_keywords:
            if keyword in text:
                return True
        return False
    
    def create_content_slide(self, prs, title, content_lines):
        """Создание слайда с содержимым"""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = self.subtitle_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.name = 'Calibri'

        # Обработка содержимого
        text_frame = content_shape.text_frame
        text_frame.clear()  # Очищаем плейсхолдер от стандартного текста и форматирования
        
        in_code_block = False
        code_lines = []
        first_paragraph = True
        
        for line in content_lines:
            if line.strip() == '[CODE_BLOCK]':
                in_code_block = True
                code_lines = []
                continue
            elif line.strip() == '[/CODE_BLOCK]':
                in_code_block = False
                # Добавляем код блок
                self.add_code_block(text_frame, code_lines)
                code_lines = []
                first_paragraph = False
                continue
            
            if in_code_block:
                code_lines.append(line)
            else:
                # Обычный текст
                if first_paragraph:
                    p = text_frame.paragraphs[0]
                    first_paragraph = False
                else:
                    p = text_frame.add_paragraph()
                
                # Проверяем, является ли строка элементом списка
                if line.strip().startswith('- '):
                    # ЭТО ЭЛЕМЕНТ СПИСКА - оставляем "- " как есть и делаем отступ
                    p.text = line.strip()  # Оставляем "- " как есть
                    p.level = 0  # Отключаем автоматические маркеры
                    
                    # Устанавливаем отступ, но НЕ отключаем маркеры для строк с "- "
                    pPr = p._element.get_or_add_pPr()
                    pPr.set('marL', '360000')  # Левый отступ
                else:
                    # ЭТО ОБЫЧНЫЙ АБЗАЦ - отключаем маркеры только для обычных абзацев
                    p.text = line
                    p.level = 0  # Отключаем автоматические маркеры
                    
                    # Полностью отключаем маркеры через XML только для обычных абзацев
                    from lxml import etree
                    pPr = p._element.get_or_add_pPr()
                    # Добавляем элемент buNone для полного отключения маркеров
                    buNone = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')

                # Применяем общие стили к тексту абзаца
                p.font.size = self.content_font_size
                p.font.color.rgb = self.content_color
                p.font.name = 'Calibri'

                # Выделение жирным ключевых слов
                if self.should_bold_text(line):
                    p.font.bold = True
    
    def add_code_block(self, text_frame, code_lines):
        """Добавление блока кода"""
        if not code_lines:
            return
            
        # Создаем параграф для кода
        p = text_frame.add_paragraph()
        p.text = '\n'.join(code_lines)
        p.font.size = self.code_font_size
        p.font.name = 'Consolas'
        p.font.color.rgb = self.code_color
        p.level = 0  # Убираем маркеры списка для кода
    
    def generate_presentations(self, content, output_dir):
        """Основная функция генерации презентаций"""
        presentations_data = self.parse_content(content)
        
        # Создаем выходную директорию если её нет
        os.makedirs(output_dir, exist_ok=True)
        
        created_files = []
        
        for i, presentation_data in enumerate(presentations_data, 1):
            # Создаем новую презентацию для каждой темы
            prs = Presentation()
            
            # Создаем титульный слайд
            title = presentation_data.get('title', f'Презентация {i}')
            subtitle_parts = []
            if presentation_data.get('level'):
                subtitle_parts.append(f"Уровень: {presentation_data['level']}")
            if presentation_data.get('module'):
                subtitle_parts.append(presentation_data['module'])
            
            subtitle = ' | '.join(subtitle_parts) if subtitle_parts else ""
            self.create_title_slide(prs, title, subtitle)
            
            # Создаем слайды с содержимым
            for slide_data in presentation_data.get('slides', []):
                slide_title = slide_data.get('title', 'Слайд')
                slide_content = slide_data.get('content', [])
                self.create_content_slide(prs, slide_title, slide_content)
            
            # Формируем имя файла
            safe_title = re.sub(r'[^\w\s-]', '', title).strip()
            safe_title = re.sub(r'[-\s]+', '_', safe_title)
            filename = f"presentation_{i:02d}_{safe_title}.pptx"
            output_path = os.path.join(output_dir, filename)
            
            # Сохраняем презентацию
            prs.save(output_path)
            created_files.append(output_path)
            print(f"Создана презентация: {output_path}")
        
        return created_files


def main():
    parser = argparse.ArgumentParser(description='Генератор презентаций из структурированного текста')
    parser.add_argument('input_file', help='Путь к входному файлу с текстом')
    parser.add_argument('-o', '--output', help='Путь к выходной директории', 
                       default='presentations')
    
    args = parser.parse_args()
    
    # Читаем входной файл
    try:
        with open(args.input_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except FileNotFoundError:
        print(f"Ошибка: Файл {args.input_file} не найден")
        return
    except Exception as e:
        print(f"Ошибка при чтении файла: {e}")
        return
    
    # Генерируем презентации
    generator = PresentationGenerator()
    created_files = generator.generate_presentations(content, args.output)
    
    print(f"\nВсего создано презентаций: {len(created_files)}")
    for file_path in created_files:
        print(f"  - {file_path}")


if __name__ == "__main__":
    main()